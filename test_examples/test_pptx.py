#!/usr/bin/env python2.7

"""
# pptx_add_image.py
# An example of adding an image to the Powerpoint slide with python_pptx
"""

import myutil
from myutil import *
import pptx
from pptx import Presentation
from pptx.util import Inches, Px
import PIL
from PIL import Image


mypath = 'files'
basewidth = 650
img = Image.open('%s/dashboard.png' % mypath)
wpercent = (basewidth/float(img.size[0]))

hsize = int((float(img.size[1])*float(wpercent)))
img = img.resize((basewidth,hsize), PIL.Image.ANTIALIAS)
img.save('%s/dashboard2.png' % mypath)

img_path = '%s/dashboard2.png' % mypath

prs = Presentation()

# --------------------------------------------------------------
blank_slidelayout = prs.slidelayouts[6]

slide = prs.slides.add_slide(blank_slidelayout)

top = Inches(1)
left = Inches(0.5)

pic = slide.shapes.add_picture(img_path, left, top)

width = Px(650)
height = int(width*wpercent)
pic = slide.shapes.add_picture(img_path, left, top, width, height)

# --------------------------------------------------------------
title_only_slidelayout = prs.slidelayouts[5]
slide = prs.slides.add_slide(title_only_slidelayout)
shapes = slide.shapes

shapes.title.text = 'Adding a Table'

cols = 2
rows = 4
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

tbl = shapes.add_table(rows, cols, left, top, width, height)

# set column widths
tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(4.0)

# write column headings
tbl.cell(0, 0).text = 'Foo'
tbl.cell(0, 1).text = 'Bar'

# write body cells
tbl.cell(1, 0).text = 'Baz'
tbl.cell(1, 1).text = 'Qux'
tbl.cell(2, 0).text = 'Baz'
tbl.cell(2, 1).text = '223'
tbl.cell(3, 0).text = '456'
tbl.cell(3, 1).text = '543'

# --------------------------------------------------------------
bullet_slidelayout = prs.slidelayouts[1]

slide = prs.slides.add_slide(bullet_slidelayout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'

tf = body_shape.textframe
tf.text = 'Find the bullet slide layout'

p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1

p = tf.add_paragraph()
p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2


prs.save('%s/test_pptx.pptx' % mypath)